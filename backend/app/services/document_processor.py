import os
from typing import List, Dict, Any
from pathlib import Path
import hashlib

# Document processors
from pypdf import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation
from bs4 import BeautifulSoup
import pandas as pd

from app.core.config import settings


class DocumentProcessor:
    """
    Handles text extraction, chunking, and processing of various document types.

    This service is responsible for reading files, extracting their text content,
    and splitting the text into smaller, manageable chunks suitable for
    embedding and vector storage.
    """
    
    def __init__(self):
        """Initializes the DocumentProcessor with chunking settings."""
        self.chunk_size = settings.chunk_size
        self.chunk_overlap = settings.chunk_overlap
        
    def extract_text(self, file_path: str, file_type: str) -> str:
        """
        Extracts text from a document based on its file type.

        This method acts as a dispatcher, calling the appropriate private
        extraction method for the given file type.

        Args:
            file_path (str): The path to the document file.
            file_type (str): The file extension (e.g., 'pdf', 'docx').

        Returns:
            str: The extracted text content from the document.

        Raises:
            ValueError: If the file type is not supported.
        """
        extractors = {
            'pdf': self._extract_pdf,
            'docx': self._extract_docx,
            'pptx': self._extract_pptx,
            'txt': self._extract_txt,
            'csv': self._extract_csv,
            'html': self._extract_html,
        }
        
        extractor = extractors.get(file_type.lower())
        if not extractor:
            raise ValueError(f"Unsupported file type: {file_type}")
            
        return extractor(file_path)
    
    def _extract_pdf(self, file_path: str) -> str:
        """
        Extracts text from a PDF file.

        Args:
            file_path (str): The path to the PDF file.

        Returns:
            str: The concatenated text from all pages.
        """
        text = ""
        with open(file_path, 'rb') as file:
            pdf_reader = PdfReader(file)
            for page in pdf_reader.pages:
                text += page.extract_text() + "\n"
        return text
    
    def _extract_docx(self, file_path: str) -> str:
        """
        Extracts text from a DOCX file.

        Args:
            file_path (str): The path to the DOCX file.

        Returns:
            str: The concatenated text from all paragraphs.
        """
        doc = DocxDocument(file_path)
        return "\n".join([paragraph.text for paragraph in doc.paragraphs])
    
    def _extract_pptx(self, file_path: str) -> str:
        """
        Extracts text from a PPTX file.

        Args:
            file_path (str): The path to the PPTX file.

        Returns:
            str: The concatenated text from all shapes on all slides.
        """
        prs = Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    
    def _extract_txt(self, file_path: str) -> str:
        """
        Extracts text from a plain text file.

        Args:
            file_path (str): The path to the TXT file.

        Returns:
            str: The content of the file.
        """
        with open(file_path, 'r', encoding='utf-8') as file:
            return file.read()
    
    def _extract_csv(self, file_path: str) -> str:
        """
        Extracts text from a CSV file by converting it to a string.

        Args:
            file_path (str): The path to the CSV file.

        Returns:
            str: A string representation of the CSV data.
        """
        df = pd.read_csv(file_path)
        return df.to_string()
    
    def _extract_html(self, file_path: str) -> str:
        """
        Extracts text from an HTML file.

        Args:
            file_path (str): The path to the HTML file.

        Returns:
            str: The visible text content of the HTML.
        """
        with open(file_path, 'r', encoding='utf-8') as file:
            soup = BeautifulSoup(file.read(), 'html.parser')
            return soup.get_text()
    
    def create_chunks(self, text: str, metadata: Dict[str, Any]) -> List[Dict[str, Any]]:
        """
        Splits a string of text into smaller, overlapping chunks.

        Each chunk is given a unique ID based on its content and includes
        the provided metadata.

        Args:
            text (str): The input text to be chunked.
            metadata (Dict[str, Any]): Metadata to be included with each chunk.

        Returns:
            List[Dict[str, Any]]: A list of chunk dictionaries, where each
                                  dictionary contains the chunk ID, content,
                                  and metadata.
        """
        words = text.split()
        chunks = []
        
        for i in range(0, len(words), self.chunk_size - self.chunk_overlap):
            chunk_words = words[i:i + self.chunk_size]
            chunk_text = " ".join(chunk_words)
            
            if chunk_text.strip():  # Skip empty chunks
                chunk_id = hashlib.md5(chunk_text.encode()).hexdigest()
                chunks.append({
                    "chunk_id": chunk_id,
                    "content": chunk_text,
                    "metadata": {
                        **metadata,
                        "chunk_index": len(chunks),
                        "start_index": i,
                        "end_index": min(i + self.chunk_size, len(words))
                    }
                })
        
        return chunks
    
    async def process_document(self, file_path: str, filename: str, file_type: str) -> List[Dict[str, Any]]:
        """
        Orchestrates the full processing of a single document.

        This method handles the entire pipeline: extracting text from the file,
        creating metadata, and splitting the text into chunks.

        Args:
            file_path (str): The path to the document file.
            filename (str): The original name of the file.
            file_type (str): The file extension.

        Returns:
            List[Dict[str, Any]]: A list of chunk dictionaries ready for ingestion.
        """
        # Extract text
        text = self.extract_text(file_path, file_type)
        
        # Create metadata
        metadata = {
            "filename": filename,
            "file_type": file_type,
            "file_path": file_path,
        }
        
        # Create chunks
        chunks = self.create_chunks(text, metadata)
        
        return chunks